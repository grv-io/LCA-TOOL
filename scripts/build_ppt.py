# -*- coding: utf-8 -*-
"""DhatuChakra — SIH Idea PPT builder.
Loads the official SIH template (unchanged headings/footers/logos) and fills it
with designed content. All numbers come from prototype/DATA_SPEC.md (verified).
Re-runnable: overwrites the output file each time.
"""
import copy
from pptx import Presentation
from pptx.util import Inches as IN, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DOC_URL = "https://github.com/grv-io/LCA-TOOL/blob/main/docs/DhatuChakra_Documentation.pdf"
REPO_URL = "https://github.com/grv-io/LCA-TOOL"
TEMPLATE = r"C:\Users\Gaurav Agrawal\Downloads\SIH-IDEA-Presentation-Format.pptx"
OUT = r"C:\Users\Gaurav Agrawal\OneDrive\Desktop\lca-tool\DhatuChakra_SIH_Idea.pptx"

# palette (DESIGN.md)
SAFFRON = RGBColor(0xE8, 0x6A, 0x17)
SAFFRON_DEEP = RGBColor(0xB8, 0x4E, 0x0A)
SAFFRON_WASH = RGBColor(0xFC, 0xEF, 0xE3)
NAVY = RGBColor(0x24, 0x46, 0x6E)
NAVY_WASH = RGBColor(0xEB, 0xF1, 0xF7)
GREEN = RGBColor(0x17, 0x80, 0x3D)
GREEN_WASH = RGBColor(0xE9, 0xF4, 0xEC)
INK = RGBColor(0x1C, 0x2B, 0x3A)
SOFT = RGBColor(0x5A, 0x6B, 0x7B)
LINE = RGBColor(0xE7, 0xE0, 0xD5)
MUSTARD_WASH = RGBColor(0xF7, 0xEF, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY_WASH = RGBColor(0xF2, 0xEE, 0xE6)

BODY = "Segoe UI"
HEAD = "Trebuchet MS"


def find(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def set_runs(para, parts, size=12, align=None, space_after=None):
    """parts = list of (text, dict(bold=, color=, italic=, font=, size=))"""
    for r in list(para.runs):
        r._r.getparent().remove(r._r)
    for text, st in parts:
        run = para.add_run()
        run.text = text
        if st.get("link"):
            run.hyperlink.address = st["link"]
        f = run.font
        f.name = st.get("font", BODY)
        f.size = Pt(st.get("size", size))
        f.bold = st.get("bold", False)
        f.italic = st.get("italic", False)
        f.color.rgb = st.get("color", INK)
    if align:
        para.alignment = align
    if space_after is not None:
        para.space_after = Pt(space_after)


def box(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, shadow=False):
    sh = slide.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = IN(0.12)
    tf.margin_right = IN(0.12)
    tf.margin_top = IN(0.07)
    tf.margin_bottom = IN(0.07)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return sh


def txt(slide, x, y, w, h):
    sh = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = sh.text_frame
    tf.word_wrap = True
    return sh


def para(tf, first=False):
    if first and not tf.paragraphs[0].runs:
        return tf.paragraphs[0]
    return tf.add_paragraph()


def bullets(tf, items, size=12.5, gap=6, color=INK, bullet="•  "):
    first = True
    for it in items:
        p = para(tf, first)
        first = False
        if isinstance(it, tuple):  # (bold-lead, rest)
            set_runs(p, [(bullet, {"color": SAFFRON_DEEP, "bold": True, "size": size}),
                         (it[0], {"bold": True, "size": size, "color": color}),
                         (it[1], {"size": size, "color": color})], space_after=gap)
        else:
            set_runs(p, [(bullet, {"color": SAFFRON_DEEP, "bold": True, "size": size}),
                         (it, {"size": size, "color": color})], space_after=gap)


def heading(tf, text, size=15, color=NAVY, first=True, after=6):
    p = para(tf, first)
    set_runs(p, [(text.upper(), {"bold": True, "size": size, "color": color, "font": HEAD})], space_after=after)
    return p


def caption_pointer(slide, text):
    """Shrink the template's prompt TextBox 8 into a small italic pointer line (kept verbatim)."""
    tb = find(slide, "TextBox 8")
    tb.left, tb.top, tb.width, tb.height = IN(0.4), IN(1.12), IN(12.5), IN(0.28)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    set_runs(tf.paragraphs[0], [(text, {"italic": True, "size": 10.5, "color": SOFT})])


def arrow(slide, x, y, w=0.3, h=0.28, direction="right", color=SOFT):
    shp = {"right": MSO_SHAPE.RIGHT_ARROW, "left": MSO_SHAPE.LEFT_ARROW,
           "up": MSO_SHAPE.UP_ARROW, "down": MSO_SHAPE.DOWN_ARROW}[direction]
    a = slide.shapes.add_shape(shp, IN(x), IN(y), IN(w), IN(h))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def diagram_box(slide, x, y, w, h, title, sub, fill, border, tcolor=INK, tsize=11, ssize=9.5):
    b = box(slide, x, y, w, h, fill=fill, line=border, line_w=1.2)
    tf = b.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_runs(tf.paragraphs[0], [(title, {"bold": True, "size": tsize, "color": tcolor})],
             align=PP_ALIGN.CENTER, space_after=1)
    if sub:
        p = tf.add_paragraph()
        set_runs(p, [(sub, {"size": ssize, "color": SOFT if tcolor == INK else tcolor})], align=PP_ALIGN.CENTER)
    return b


prs = Presentation(TEMPLATE)
s1, s2, s3, s4, s5, s6, s7 = list(prs.slides)

# ============================================================ SLIDE 1 — TITLE
tb = find(s1, "TextBox 9")
tf = tb.text_frame
tf.clear()
tf.word_wrap = True
rows = [
    ("Problem Statement ID  ", "48 (internal round)"),
    ("Problem Statement Title  ", "AI-Driven Life Cycle Assessment (LCA) Tool for Advancing Circularity and Sustainability in Metallurgy and Mining"),
    ("Theme  ", "Miscellaneous (Sustainability)"),
    ("PS Category  ", "Software"),
    ("Team ID  ", "________________"),
    ("Team Name  ", "________________"),
]
first = True
for label, val in rows:
    p = para(tf, first)
    first = False
    set_runs(p, [("‣ ", {"color": SAFFRON, "bold": True, "size": 15}),
                 (label, {"bold": True, "size": 15, "color": INK}),
                 (val, {"size": 14.5, "color": NAVY})], space_after=10)
p = tf.add_paragraph()
set_runs(p, [("", {"size": 6})], space_after=2)
p = tf.add_paragraph()
set_runs(p, [("Idea:  ", {"bold": True, "size": 15, "color": INK}),
             ("DhatuChakra ", {"bold": True, "size": 16, "color": SAFFRON_DEEP, "font": HEAD}),
             ("(धातुचक्र)", {"size": 14, "color": SOFT, "font": "Nirmala UI"}),
             ("  —  हर टन का हिसाब · every tonne, accounted for", {"italic": True, "size": 12.5, "color": SOFT, "font": "Nirmala UI"})])

# ============================================================ SLIDE 2 — IDEA
caption_pointer(s2, "Proposed Solution (Describe your Idea/Solution/Prototype)  ·  How it addresses the problem  ·  Innovation and uniqueness of the solution")

b = box(s2, 0.4, 1.42, 7.55, 1.02, fill=SAFFRON_WASH, line=SAFFRON, line_w=1.2)
tf = b.text_frame
set_runs(tf.paragraphs[0], [("DhatuChakra ", {"bold": True, "size": 22, "color": SAFFRON_DEEP, "font": HEAD}),
                            ("(धातुचक्र)", {"size": 16, "color": INK, "font": "Nirmala UI"})], space_after=1)
p = tf.add_paragraph()
set_runs(p, [("AI-assisted life-cycle assessment & circularity engine for Indian metals — 5 inputs in, ISO-14044-style answers out.",
              {"size": 11.5, "color": INK})])

b = box(s2, 8.15, 1.42, 4.78, 1.02, fill=NAVY_WASH, line=NAVY, line_w=1.0)
tf = b.text_frame
set_runs(tf.paragraphs[0], [("Working prototype — see it run", {"bold": True, "size": 11, "color": NAVY})], space_after=2)
_links2 = [
    ("Demo video:  ", "<add link>", None),
    ("Documentation:  ", "LCA-TOOL/docs · technical PDF", DOC_URL),
    ("Code + prototype:  ", "github.com/grv-io/LCA-TOOL", REPO_URL),
]
for lead, val, url in _links2:
    p = tf.add_paragraph()
    st2 = {"italic": url is None, "size": 10.5, "color": SAFFRON_DEEP if url else SOFT}
    if url: st2["link"] = url
    set_runs(p, [("▸ ", {"color": SAFFRON_DEEP, "bold": True, "size": 10.5}),
                 (lead, {"bold": True, "size": 10.5, "color": INK}),
                 (val, st2)], space_after=1)

b = box(s2, 0.4, 2.62, 6.1, 4.12, fill=WHITE, line=LINE, line_w=1)
tf = b.text_frame
heading(tf, "The problem → our solution")
bullets(tf, [
    ("LCA today is out of reach: ", "25+ process parameters, database licences (ecoinvent ≈ €4,000/yr), weeks of consultants — most Indian plants simply skip it"),
    ("DhatuChakra needs 5–8 inputs ", "— or one plain sentence (\"aluminium smelter in Odisha, 40% scrap\"); an LLM parses it into a structured scenario"),
    ("AI fills every missing parameter ", "with a confidence score — estimates are visibly flagged, never hidden"),
    ("Full picture per tonne: ", "GWP, energy, water, acidification + Material Circularity Indicator (Ellen MacArthur v3)"),
    ("Decision tool, not a calculator: ", "linear-vs-circular compare with live sliders + auditable PDF report"),
], size=11.5, gap=7)

b = box(s2, 6.8, 2.62, 6.13, 4.12, fill=WHITE, line=LINE, line_w=1)
tf = b.text_frame
heading(tf, "Innovation & uniqueness")
bullets(tf, [
    ("Not a black box: ", "every figure traces to a cited public factor — source, year, region shown in-app"),
    ("India-first: ", "clickable state-grid map (CEA-derived: Himachal 0.18 → Chhattisgarh 0.92 kg CO₂/kWh) + CBAM export cost in ₹ crore"),
    ("Uncertainty built in: ", "real 1,000-run Monte Carlo in the prototype — p05–p95 range + histogram on every result"),
    ("Instant what-if: ", "surrogate model answers sliders in milliseconds — prototype already proves the UX"),
    ("₹0 data cost: ", "built entirely on free public inventories (IAI, worldsteel, EF 3.1, IDEMAT, US LCI)"),
], size=11.5, gap=6)
p = tf.add_paragraph()
set_runs(p, [("Verified: recycled aluminium route = 16.6 → 0.92 t CO₂e/t (−94%), inside IAI published ranges.",
              {"bold": True, "size": 11, "color": GREEN})])

# ============================================================ SLIDE 3 — TECHNICAL
caption_pointer(s3, "Technologies to be used (programming languages, frameworks, hardware)  ·  Methodology and process for implementation (Flow Charts / Images / working prototype)")

b = box(s3, 0.4, 1.48, 3.6, 5.3, fill=WHITE, line=LINE, line_w=1)
tf = b.text_frame
heading(tf, "Stack")
stack = [
    ("Frontend  ", "Next.js · TypeScript · Tailwind · Recharts · Plotly (Sankey)"),
    ("Backend  ", "Python FastAPI · Pydantic · PostgreSQL"),
    ("LCA engine  ", "factor-matrix calc · EMF MCI v3 · Brightway2.5-ready"),
    ("AI / ML  ", "LightGBM quantile imputer (p10–p90 confidence) · surrogate regressor R² > 0.95, < 20 ms · NumPy Monte Carlo (1,000 runs)"),
    ("LLM  ", "Claude API — plain text → scenario JSON + grounded recommendations"),
    ("Reports  ", "Jinja2 → WeasyPrint PDF"),
    ("Deploy  ", "Docker · Render · Vercel"),
]
firstp = True
for lead, rest in stack:
    p = para(tf, firstp)
    firstp = False
    set_runs(p, [(lead, {"bold": True, "size": 10.5, "color": NAVY}),
                 (rest, {"size": 10.5, "color": INK})], space_after=6)
p = tf.add_paragraph()
set_runs(p, [("Prototype today: ", {"bold": True, "size": 10.5, "color": SAFFRON_DEEP}),
             ("client-side HTML/JS, offline — 3 metals, India state map, k-NN imputer (500 scenarios), 1,000-run Monte Carlo, CBAM, Hindi UI.", {"size": 10.5, "color": INK})])

# ---- flow diagram
t = txt(s3, 4.25, 1.42, 8.6, 0.35)
set_runs(t.text_frame.paragraphs[0], [("How a query flows", {"bold": True, "size": 13, "color": NAVY, "font": HEAD})])

y1, h1 = 1.95, 0.95
diagram_box(s3, 4.25, y1, 1.85, h1, "User input", "5 fields — or one plain-English sentence", NAVY_WASH, NAVY)
arrow(s3, 6.12, y1 + 0.33)
diagram_box(s3, 6.45, y1, 1.85, h1, "LLM parser", "text → structured scenario JSON", SAFFRON_WASH, SAFFRON)
arrow(s3, 8.32, y1 + 0.33)
diagram_box(s3, 8.65, y1, 1.85, h1, "AI imputer", "fills 25+ params + confidence score", SAFFRON_WASH, SAFFRON)
arrow(s3, 10.52, y1 + 0.33)
diagram_box(s3, 10.85, y1, 2.0, h1, "LCA engine", "quantities × emission factors", NAVY, NAVY, tcolor=WHITE)

tag = box(s3, 6.45, 1.66, 4.05, 0.24, fill=None, line=SAFFRON, line_w=0.75)
tf = tag.text_frame
tf.margin_top = IN(0.0)
tf.margin_bottom = IN(0.0)
set_runs(tf.paragraphs[0], [("AI layer", {"size": 8.5, "color": SAFFRON_DEEP, "bold": True})], align=PP_ALIGN.CENTER)

arrow(s3, 11.0, y1 + h1 + 0.03, w=0.26, h=0.40, direction="up", color=SOFT)
db = box(s3, 4.25, 3.38, 6.9, 0.72, fill=GREY_WASH, line=SOFT, line_w=0.9, shape=MSO_SHAPE.CAN)
tf = db.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_runs(tf.paragraphs[0], [("Factor library (India-first): ", {"bold": True, "size": 10.5, "color": INK}),
                            ("IAI · worldsteel · CEA CO₂ Baseline · EF 3.1 · IDEMAT — every factor cited (source · year · region)",
                             {"size": 10, "color": SOFT})], align=PP_ALIGN.CENTER)
arrow(s3, 11.72, y1 + h1 + 0.03, w=0.26, h=1.60, direction="down", color=SOFT)

y2, h2 = 4.55, 1.05
diagram_box(s3, 10.4, y2, 2.45, h2, "Monte Carlo", "1,000 runs → p05–p95 ranges + histogram", MUSTARD_WASH, RGBColor(0xC9, 0xA2, 0x27))
arrow(s3, 10.06, y2 + 0.38, direction="left")
diagram_box(s3, 7.38, y2, 2.6, h2, "Results / tonne", "GWP · energy · water · acidification + MCI chakra", GREEN_WASH, GREEN)
arrow(s3, 7.04, y2 + 0.38, direction="left")
diagram_box(s3, 4.25, y2, 2.7, h2, "Dashboard & report", "Sankey · compare sliders · recommendations · PDF", WHITE, NAVY)

t = txt(s3, 4.25, 5.85, 8.65, 0.9)
tf = t.text_frame
set_runs(tf.paragraphs[0], [("Working prototype (this exact methodology, client-side):  ", {"bold": True, "size": 11, "color": INK}),
                            ("demo video — <add link>   ·   ", {"italic": True, "size": 11, "color": SAFFRON_DEEP}),
                            ("technical documentation (PDF in repo)", {"size": 11, "color": SAFFRON_DEEP, "link": DOC_URL})])

# ============================================================ SLIDE 4 — FEASIBILITY
caption_pointer(s4, "Analysis of the feasibility of the idea  ·  Potential challenges and risks  ·  Strategies for overcoming these challenges")

b = box(s4, 0.4, 1.48, 5.0, 5.3, fill=WHITE, line=LINE, line_w=1)
tf = b.text_frame
heading(tf, "Why this is buildable")
bullets(tf, [
    ("100% free data — ", "IAI, worldsteel, CEA, EF 3.1, IDEMAT, US LCI: no paid-licence blocker anywhere in the pipeline"),
    ("Already validated: ", "prototype outputs 16.6 / 0.92 / 2.19 / 0.70 t CO₂e/t (Al primary · Al recycled · BF-BOF · EAF) sit inside published ranges — IAI 16–20 & 0.5–1.5, worldsteel 1.9–2.4 & 0.4–1.0"),
    ("Working prototype live ", "— 3 metals, state-grid map, k-NN imputer, Monte Carlo, CBAM calculator, Hindi UI; 4-week full-build plan written"),
    ("ML needs no scarce dataset — ", "imputer & surrogate train on physics-generated scenarios from our own engine"),
], size=11.5, gap=8)

tbl_x, tbl_y, tbl_w = 5.7, 1.48, 7.2
rows_data = [
    ("Challenge / Risk", "Strategy"),
    ("LCI data gaps for Indian plants", "Curated 100+ factor table, every row cited (source · year · region); fallback chain IN → world average"),
    ("Wrong numbers = zero credibility", "Validation suite vs IAI / worldsteel published values runs on every change; < 10% error target"),
    ("\u201CAI is a black box\u201D objection", "Confidence score on every estimate + provenance table + uncertainty ranges on screen"),
    ("Users can't fill 25+ parameters", "Plain-English entry + AI imputation — only 5–8 inputs needed"),
    ("Scope creep in hackathon time", "v1 locked: Al + steel, 4 impact categories; copper & Brightway as stretch"),
]
tbl = s4.shapes.add_table(len(rows_data), 2, IN(tbl_x), IN(tbl_y), IN(tbl_w), IN(5.3)).table
tbl.columns[0].width = IN(2.55)
tbl.columns[1].width = IN(4.65)
for ri, (c1, c2) in enumerate(rows_data):
    for ci, val in enumerate((c1, c2)):
        cell = tbl.cell(ri, ci)
        cell.margin_left = IN(0.1)
        cell.margin_right = IN(0.1)
        cell.margin_top = IN(0.05)
        cell.margin_bottom = IN(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfc = cell.text_frame
        tfc.word_wrap = True
        if ri == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            set_runs(tfc.paragraphs[0], [(val, {"bold": True, "size": 12, "color": WHITE, "font": HEAD})])
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else SAFFRON_WASH if False else WHITE
            set_runs(tfc.paragraphs[0], [(val, {"size": 10.5, "color": INK, "bold": ci == 0})])

# ============================================================ SLIDE 5 — IMPACT
caption_pointer(s5, "Potential impact on the target audience  ·  Benefits of the solution (social, economic, environmental, etc.)")

stats = [
    ("−94%", "CO₂ on the recycled aluminium route — 16.6 → 0.92 t CO₂e per tonne", GREEN, GREEN_WASH),
    ("≈ 4.9 t", "CO₂e/t removed by a 50% renewable PPA on a primary Al smelter (India grid)", NAVY, NAVY_WASH),
    ("₹0 · minutes", "database cost & assessment time — vs €4k/yr licences and weeks of consulting", SAFFRON_DEEP, SAFFRON_WASH),
]
sx = 0.4
for big, small, col, wash in stats:
    b = box(s5, sx, 1.48, 4.1, 1.42, fill=wash, line=col, line_w=1.2)
    tf = b.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_runs(tf.paragraphs[0], [(big, {"bold": True, "size": 24, "color": col, "font": HEAD})], align=PP_ALIGN.CENTER, space_after=2)
    p = tf.add_paragraph()
    set_runs(p, [(small, {"size": 10.5, "color": INK})], align=PP_ALIGN.CENTER)
    sx += 4.27

b = box(s5, 0.4, 3.12, 6.1, 3.65, fill=WHITE, line=LINE, line_w=1)
tf = b.text_frame
heading(tf, "Who gains")
bullets(tf, [
    ("Metal producers & MSMEs — ", "LCA answers in minutes; decide scrap %, energy sourcing and routes with numbers, not guesses"),
    ("Ministry of Mines & regulators — ", "standardised, comparable, auditable plant data; supports the National Critical Mineral Mission and circular-economy policy"),
    ("Exporters — ", "EU CBAM's definitive regime starts 2026: embedded-emissions reporting for steel & aluminium exports becomes mandatory — exactly what this tool produces"),
], size=11.5, gap=8)

b = box(s5, 6.8, 3.12, 6.13, 3.65, fill=WHITE, line=LINE, line_w=1)
tf = b.text_frame
heading(tf, "Benefits")
bullets(tf, [
    ("Environmental — ", "quantifies the shift to circular routes; pinpoints hotspot stages (smelting ≈ 62% of primary-Al footprint)"),
    ("Economic — ", "saves licence + consultant costs; higher circularity = lower import dependence on ore & energy"),
    ("Social — ", "small plants get the same analysis capability as corporates; methodology fully transparent"),
    ("Strategic — ", "an India-first open emission-factor library becomes reusable national infrastructure"),
], size=11.5, gap=7)

# ============================================================ SLIDE 6 — REFERENCES
caption_pointer(s6, "Details / Links of the reference and research work")

b = box(s6, 0.4, 1.48, 7.9, 5.3, fill=WHITE, line=LINE, line_w=1)
tf = b.text_frame
heading(tf, "Research & data sources")
refs = [
    ("ISO 14040 / 14044", " — LCA principles, requirements & guidelines · iso.org"),
    ("Ellen MacArthur Foundation", " — Material Circularity Indicator methodology v3 · ellenmacarthurfoundation.org"),
    ("International Aluminium Institute", " — LCI data & GHG pathways · international-aluminium.org"),
    ("worldsteel", " — Life Cycle Inventory methodology report · worldsteel.org"),
    ("CEA CO₂ Baseline Database v19", " — India grid emission factor · cea.nic.in"),
    ("US LCI (NREL) / EU EF 3.1 · ELCD", " — lcacommons.gov · eplca.jrc.ec.europa.eu"),
    ("IDEMAT 2024, TU Delft", " — open LCI factors · idematapp.com"),
    ("IBM Indian Minerals Yearbook", " — Ministry of Mines · ibm.gov.in"),
    ("EU CBAM Regulation 2023/956", " — carbon reporting for metal imports · taxation-customs.ec.europa.eu"),
    ("NITI Aayog", " — circular-economy strategy papers for metals · niti.gov.in"),
]
firstp = True
for i, (lead, rest) in enumerate(refs, 1):
    p = para(tf, firstp)
    firstp = False
    set_runs(p, [(f"{i:>2}.  ", {"bold": True, "size": 11, "color": SAFFRON_DEEP}),
                 (lead, {"bold": True, "size": 11, "color": INK}),
                 (rest, {"size": 11, "color": SOFT})], space_after=5)

b = box(s6, 8.5, 1.48, 4.43, 2.5, fill=NAVY_WASH, line=NAVY, line_w=1.0)
tf = b.text_frame
heading(tf, "Our links", size=13)
_links6 = [
    ("Demo video", "<add link>", None),
    ("Documentation (PDF)", "LCA-TOOL/docs/DhatuChakra_Documentation.pdf", DOC_URL),
    ("Code + prototype", "github.com/grv-io/LCA-TOOL", REPO_URL),
]
for lbl, val, url in _links6:
    p = tf.add_paragraph()
    st6 = {"italic": url is None, "size": 11, "color": SAFFRON_DEEP if url else SOFT}
    if url: st6["link"] = url
    set_runs(p, [("▸ ", {"color": SAFFRON_DEEP, "bold": True, "size": 11}),
                 (lbl + ":  ", {"bold": True, "size": 11, "color": INK}),
                 (val, st6)], space_after=5)

b = box(s6, 8.5, 4.2, 4.43, 2.58, fill=GREEN_WASH, line=GREEN, line_w=1.0)
tf = b.text_frame
heading(tf, "Reproducible by design", size=13, color=GREEN)
p = tf.add_paragraph()
set_runs(p, [("Every number in this deck — 16.6, 0.92, 2.19, 0.70 t CO₂e/t, MCI 0.16→0.83, 62% electricity share — "
              "is computed from the sources listed; formulas and validation are in the linked technical PDF.",
              {"size": 11, "color": INK})])

# ============================================================ delete instructions slide 7
xml_slides = prs.slides._sldIdLst
xml_slides.remove(list(xml_slides)[6])

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(prs.slides._sldIdLst))
